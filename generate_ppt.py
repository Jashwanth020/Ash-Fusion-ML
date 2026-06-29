from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
# Use a blank slide layout
blank_slide_layout = prs.slide_layouts[5] # Title only
slide = prs.slides.add_slide(blank_slide_layout)

# Add Title
title = slide.shapes.title
title.text = "Ash Fusion Temperature Predictions (CLSO Reconciled)"

# Add Table
rows = 6
cols = 5
left = Inches(0.5)
top = Inches(1.5)
width = Inches(8.5)
height = Inches(2.0)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Headers
headers = ['Blend Sample', 'DT (°C)', 'ST (°C)', 'HT (°C)', 'FT (°C)']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True

# Data
data = [
    ['100 coal', '1168.1', '1244.5', '1254.5', '1351.2'],
    ['5% BC', '1151.1', '1239.3', '1249.3', '1292.7'],
    ['10% BC', '1145.9', '1218.3', '1228.3', '1303.2'],
    ['15% BC', '1162.7', '1260.1', '1270.1', '1331.9'],
    ['100% BC', '1117.1', '1156.9', '1192.3', '1252.5']
]

for row_idx, row_data in enumerate(data):
    for col_idx, cell_data in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_data

# Formatting table font size
for row in table.rows:
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)

# Add Interpretation text box
left_txt = Inches(0.5)
top_txt = Inches(4.0)
width_txt = Inches(8.5)
height_txt = Inches(3.0)

txBox = slide.shapes.add_textbox(left_txt, top_txt, width_txt, height_txt)
tf = txBox.text_frame
tf.word_wrap = True

p1 = tf.add_paragraph()
p1.text = "Key Interpretations:"
p1.font.bold = True
p1.font.size = Pt(18)

bullets = [
    "Alkaline Fluxing Effect: Increasing the biomass proportion generally lowers the ash fusion temperatures across all stages (compare 100% coal to 100% BC). This is driven by the introduction of alkaline species (K₂O, CaO) which break down the aluminosilicate network.",
    "Physical Validity (CLSO): Raw predictions contained overlaps (ST > HT). The Constrained Least Squares Optimization (CLSO) algorithm successfully forced the predictions to maintain the strict thermodynamic order (DT ≤ ST ≤ HT ≤ FT).",
    "Non-linear Blending Behavior: The 15% BC blend shows a localized peak in fusion temperatures compared to 10% BC, highlighting that complex quaternary/quinary oxide interactions can create eutectic fluctuations before dropping drastically at pure biomass."
]

for b in bullets:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(14)
    p.level = 0

prs.save(r'P:\MVR\new_data_exp\AFT_Predictions_Slide.pptx')
print("PPT successfully generated.")
